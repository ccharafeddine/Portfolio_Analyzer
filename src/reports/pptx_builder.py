"""Client-facing PowerPoint report.

Builds a polished, dark-themed 16:9 deck covering the full analysis: an executive
summary, growth/performance, risk, optimization, attribution, income, forecast,
retirement plan, and tax — one section per slide, each with a rendered chart and
key figures. Charts are exported to PNG via kaleido.

``build_pptx`` is pure (no Qt); the caller passes an optional logo PNG.
"""

from __future__ import annotations

from io import BytesIO
from typing import Optional

import numpy as np

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.charts import plotly_charts as charts

# Palette (matches the app's Institutional theme).
BG = RGBColor(0x0B, 0x11, 0x20)
CARD = RGBColor(0x15, 0x1D, 0x2E)
TEXT = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _fmt_pct(v, d=1):
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "—"
    return f"{v * 100:.{d}f}%"


def _fmt_ratio(v, d=2):
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "—"
    return f"{v:.{d}f}"


def _fmt_money(v):
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "—"
    if abs(v) >= 1e6:
        return f"${v / 1e6:,.2f}M"
    return f"${v:,.0f}"


class _Deck:
    def __init__(self, logo_png: Optional[bytes]):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank = self.prs.slide_layouts[6]
        self.logo_png = logo_png

    def slide(self):
        s = self.prs.slides.add_slide(self._blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = BG
        return s

    def text(self, slide, text, left, top, width, height, size, color=TEXT,
             bold=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for line in str(text).split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            r = p.add_run()
            r.text = line
            f = r.font
            f.size = Pt(size)
            f.bold = bold
            f.color.rgb = color
            f.name = "Segoe UI"
        return tb

    def header(self, slide, title):
        self.text(slide, title, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8),
                  28, bold=True)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.15),
                                     Inches(1.6), Pt(4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

    def chart(self, slide, fig, left, top, width, height):
        try:
            fig.update_layout(width=1100, height=560)
            png = fig.to_image(format="png", scale=2)
            slide.shapes.add_picture(BytesIO(png), left, top, width=width, height=height)
        except Exception:
            self.text(slide, "(chart unavailable)", left, top, width, Inches(0.5), 12, MUTED)

    def metric_cards(self, slide, items, top):
        """A row of up to 6 metric cards: (label, value, color)."""
        n = len(items)
        gap = Inches(0.2)
        total_w = SLIDE_W - Inches(1.2)
        card_w = (total_w - gap * (n - 1)) / n
        left = Inches(0.6)
        for label, value, color in items:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                          card_w, Inches(1.15))
            card.fill.solid()
            card.fill.fore_color.rgb = CARD
            card.line.color.rgb = RGBColor(0x2D, 0x3A, 0x50)
            card.line.width = Pt(0.75)
            card.shadow.inherit = False
            self.text(slide, label.upper(), left + Inches(0.15), top + Inches(0.12),
                      card_w - Inches(0.3), Inches(0.3), 10, MUTED, bold=True)
            self.text(slide, value, left + Inches(0.15), top + Inches(0.42),
                      card_w - Inches(0.3), Inches(0.6), 22, color, bold=True)
            left += card_w + gap

    def bullets(self, slide, text, left, top, width, height, size=14):
        self.text(slide, text, left, top, width, height, size, MUTED)


def _headline_metrics(results):
    a = results.active
    ret = a.ann_return if a else None
    excess = None
    if a and results.passive:
        excess = a.ann_return - results.passive.ann_return
    alpha = beta = None
    if results.capm_results:
        alpha = float(np.mean([r.alpha for r in results.capm_results])) * 12
        beta = float(np.mean([r.beta for r in results.capm_results]))
    return [
        ("Total Return", _fmt_pct(ret), GREEN if (ret or 0) >= 0 else RED),
        ("Sharpe", _fmt_ratio(a.sharpe) if a else "—", TEXT),
        ("Max Drawdown", _fmt_pct(a.max_dd) if a else "—", RED),
        ("Volatility", _fmt_pct(a.ann_vol) if a else "—", TEXT),
        ("Alpha", _fmt_pct(alpha) if alpha is not None else "—",
         GREEN if (alpha or 0) >= 0 else RED),
        ("Beta", _fmt_ratio(beta) if beta is not None else "—", TEXT),
    ]


def build_pptx(results, path: str, logo_png: Optional[bytes] = None) -> None:
    d = _Deck(logo_png)
    cfg = results.config
    interp = results.interpretations or {}

    # ── Title slide ──
    s = d.slide()
    if logo_png:
        try:
            s.shapes.add_picture(BytesIO(logo_png), Inches(0.6), Inches(0.55),
                                 width=Inches(1.2), height=Inches(1.2))
        except Exception:
            pass
    d.text(s, "Portfolio Analysis", Inches(0.6), Inches(2.6), Inches(12), Inches(1.2),
           48, bold=True)
    universe = ", ".join(cfg.tickers[:10]) + (" …" if len(cfg.tickers) > 10 else "")
    d.text(s, f"{universe}", Inches(0.62), Inches(3.7), Inches(12), Inches(0.6), 18, ACCENT)
    _cash = float(getattr(cfg, "cash", 0.0) or 0.0)
    if _cash > 0:
        _cap_str = (
            f"{_fmt_money(cfg.capital)} total "
            f"({_fmt_money(cfg.capital - _cash)} invested + {_fmt_money(_cash)} cash)"
        )
    else:
        _cap_str = f"{_fmt_money(cfg.capital)} invested"
    d.text(s,
           f"{cfg.start_str} to {cfg.end_str}    •    Benchmark {cfg.benchmark}    •    "
           f"{_cap_str}",
           Inches(0.62), Inches(4.3), Inches(12), Inches(0.6), 14, MUTED)

    # ── Executive summary ──
    s = d.slide()
    d.header(s, "Executive Summary")
    d.metric_cards(s, _headline_metrics(results), Inches(1.5))
    if interp.get("performance") or interp.get("executive_summary"):
        d.bullets(s, interp.get("executive_summary") or interp.get("performance"),
                  Inches(0.6), Inches(3.0), Inches(12.1), Inches(4), 15)

    # ── Chart sections (conditional) ──
    def chart_slide(title, fig, note=None):
        s = d.slide()
        d.header(s, title)
        d.chart(s, fig, Inches(0.6), Inches(1.4), Inches(8.4), Inches(4.3))
        if note:
            d.bullets(s, note, Inches(9.2), Inches(1.6), Inches(3.5), Inches(5), 13)

    try:
        growth = {}
        for k, ps in (("Active", results.active), ("Passive", results.passive), ("ORP", results.orp)):
            if ps is not None:
                growth[k] = ps.values
        if growth:
            chart_slide("Growth of Capital", charts.growth_chart(growth, cfg.capital),
                        interp.get("performance"))
    except Exception:
        pass

    try:
        # Active portfolio allocation — holdings + a Cash slice when cash is held.
        cap = float(cfg.capital or 0.0)
        cash = float(getattr(cfg, "cash", 0.0) or 0.0)
        alloc = {t: float(w) for t, w in (cfg.weights or {}).items()}
        if cash > 0 and cap > 0:
            y = max(cap - cash, 0.0) / cap
            alloc = {t: w * y for t, w in alloc.items()}
            alloc["Cash"] = cash / cap
        if alloc:
            chart_slide("Portfolio Allocation",
                        charts.allocation_donut(alloc, "Active Portfolio Allocation"))
    except Exception:
        pass

    try:
        if results.active and results.passive:
            chart_slide("Performance vs Benchmark",
                        charts.outperformance_chart(results.active.values, results.passive.values))
    except Exception:
        pass

    try:
        dd = {}
        for k, ps in (("Active", results.active), ("Passive", results.passive)):
            if ps is not None:
                dd[k] = ps.values
        if dd:
            chart_slide("Risk: Drawdowns", charts.drawdown_chart(dd), interp.get("risk"))
    except Exception:
        pass

    try:
        orp = results.orp_optimization
        if orp is not None:
            rets_m = results.monthly_returns
            cols = [t for t in cfg.tickers if t in rets_m.columns and t != cfg.benchmark]
            mu = ((1 + rets_m[cols].mean()) ** 12 - 1) if cols else None
            vol = (rets_m[cols].std() * np.sqrt(12)) if cols else None
            fig = charts.efficient_frontier_chart(
                orp.frontier_vols, orp.frontier_returns, orp.expected_vol,
                orp.expected_return, cfg.risk_free_rate, asset_vols=vol, asset_returns=mu)
            chart_slide("Optimization: Efficient Frontier", fig, interp.get("optimization"))
            chart_slide("Optimal Portfolio Weights",
                        charts.weights_bar(orp.weights, "ORP Weights (Max-Sharpe)"))
    except Exception:
        pass

    try:
        if results.asset_attribution is not None and not results.asset_attribution.empty:
            chart_slide("Attribution",
                        charts.attribution_chart(results.asset_attribution,
                                                 title="Brinson–Fachler Attribution"),
                        interp.get("capm"))
    except Exception:
        pass

    try:
        if results.income_summary is not None and not results.income_summary.empty:
            chart_slide("Dividend Income",
                        charts.income_bar_chart(results.income_summary), interp.get("income"))
    except Exception:
        pass

    try:
        if results.simulations:
            sim = results.simulations[0]
            fig = charts.simulation_fan_chart(
                starting_value=sim.starting_value, paths=sim.paths,
                horizon_days=sim.horizon_days, method_name=sim.name)
            chart_slide("Forecast: Monte Carlo", fig, interp.get("simulation"))
    except Exception:
        pass

    try:
        plan = results.plan_result
        if plan is not None:
            fig = charts.simulation_fan_chart(
                starting_value=plan.starting_value, paths=plan.paths,
                horizon_days=plan.horizon_days, method_name=f"{plan.horizon_years}-Year Projection")
            note = (
                f"Success probability: {plan.success_prob:.0%}\n"
                f"Depletion risk: {plan.depletion_prob:.0%}\n"
                f"Median ending: {_fmt_money(plan.median_terminal)}\n"
            )
            if plan.safe_withdrawal_rate is not None:
                note += f"Safe withdrawal rate: {plan.safe_withdrawal_rate:.1%}\n"
            chart_slide("Retirement Plan", fig, note)
    except Exception:
        pass

    try:
        tm = results.tax_metrics
        if tm:
            s = d.slide()
            d.header(s, "Tax Analysis")
            d.metric_cards(s, [
                ("Unrealized Gain", _fmt_money(tm.get("unrealized_gain", 0)), TEXT),
                ("Harvestable Losses", _fmt_money(tm.get("harvest_potential", 0)), GREEN),
                ("Loss Candidates", str(tm.get("n_harvest", 0)), TEXT),
                ("Est. Tax", _fmt_money(tm.get("estimated_tax", 0)), RED),
            ], Inches(1.5))
    except Exception:
        pass

    # ── Disclosures ──
    s = d.slide()
    d.header(s, "Important Disclosures")
    d.bullets(
        s,
        "This report is generated by Portfolio Analyzer for informational and educational "
        "purposes only. It is not investment, tax, or financial advice, and is not an offer or "
        "solicitation to buy or sell any security. Past performance and backtested or simulated "
        "results do not guarantee future results; simulated results have inherent limitations. "
        "Figures rely on third-party market data that may contain errors. Consult a qualified "
        "professional before making investment decisions.",
        Inches(0.6), Inches(1.5), Inches(12.1), Inches(5), 13)

    d.prs.save(path)
